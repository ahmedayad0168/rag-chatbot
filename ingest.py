# modules
import os
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import pandas as pd
import pytesseract
from PIL import Image
from sentence_transformers import SentenceTransformer
import database
from config import MODEL_NAME, CHUNK_SIZE

database.init_db()

model = SentenceTransformer(MODEL_NAME)

# read .pdf
def read_pdf(file_path):
    reader = PdfReader(file_path)
    text = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:  # make sure it's not None
            text.append(page_text.strip())
    return " ".join(text)

# read .pptx
def read_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:  # check if shape contains text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text.append(run.text.strip())
    return " ".join(text).strip()

# read .txt
def read_txt(file_path):
    text = []
    with open(file_path, "r", encoding="utf-8") as file:
        for line in file:
            if line.strip():  # skip empty lines
                text.append(line.strip())
    return " ".join(text)

# read .docx
def read_docx(file_path):
    doc = docx.Document(file_path)
    text = []
    for para in doc.paragraphs:
        if para.text.strip():  # skip empty paragraphs
            text.append(para.text.strip())
    return " ".join(text)

# read excel_file .csv or xlsx
def read_excel(file_path):
    if file_path.endswith(".csv"):
        df = pd.read_csv(file_path)
    elif file_path.endswith((".xlsx", ".xls")):
        df = pd.read_excel(file_path)
    else:
        raise ValueError("Unsupported file format. Use .csv, .xlsx, or .xls")

    text = []
    for _, row in df.iterrows():
        # Convert all values in the row to strings and join them
        row_text = " ".join(str(value).strip() for value in row if pd.notna(value))
        text.append(row_text)
    
    return " ".join(text)

# read image
def read_image(file_path, lang="eng"):
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img, lang=lang)
    return " ".join(text.split()) # clean up extra whitespace and return

# chunks
def chunk_text(text, chunk_size=150):
    """Split text into chunks of up to `chunk_size` words (not characters)."""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunk = " ".join(words[i:i + chunk_size])
        chunks.append(chunk)
    return chunks

# ingest file
def ingest_file(file_path, chunk_size):
    """
    Read file at file_path, create chunks, embed and save to DB.
    IMPORTANT: this function DOES NOT copy the file - it expects the file is already at file_path.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(file_path)

    ext = os.path.splitext(file_path)[1].lower()
    readers = {
        ".pdf": read_pdf,
        ".pptx": read_pptx,
        ".txt": read_txt,
        ".docx": read_docx,
        ".csv": read_excel,
        ".xlsx": read_excel,
        ".xls": read_excel,
    }

    if ext in readers:
        text = readers[ext](file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        text = read_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    chunks = chunk_text(text, chunk_size)
    if not chunks:
        return f"No text found in {file_path}"

    conn = database.get_connection()
    cur = conn.cursor()
    cur.execute("INSERT OR IGNORE INTO documents (file_path) VALUES (?)", (file_path,))
    conn.commit()
    cur.execute("SELECT id FROM documents WHERE file_path=?", (file_path,))
    document_id = cur.fetchone()[0]
    conn.close()

    vectors = model.encode(chunks, batch_size=32, show_progress_bar=True)
    database.save_chunks(document_id, chunks, vectors)
    
    # i = 1
    # for chunk in chunks:
    #     print(f"chunk {i}: ")
    #     print(chunk)
    #     i += 1

    return f"{file_path} ingested with {len(chunks)} chunks."

# main
if __name__ == "__main__":
    sample = r"uploads/Session 1.pptx"
    print(ingest_file(sample, chunk_size=150))
